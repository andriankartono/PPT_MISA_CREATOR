from bs4 import BeautifulSoup
import requests
from pptx import Presentation
from pptx.util import Pt
import tkinter as tk
from tkinter import ttk, filedialog
import datetime

#Get the date of the closest sunday to use as default value
today= datetime.date.today()
while today.weekday() != 6:
    today += datetime.timedelta(1)

#GUI Input
win= tk.Tk()

#Variables that will be defined by user Input
tanggal=0
bulan=0
tahun=0
word_limit = 500
fontsize=24

directory = "D:\\"

def get_text():
   global entry_tanggal, entry_bulan, entry_tahun, directory, entry_fontsize, entry_wordcount
   global tanggal, bulan, tahun, word_limit, fontsize
   tanggal= entry_tanggal.get()
   bulan= entry_bulan.get()
   tahun= entry_tahun.get()
   fontsize= int(entry_fontsize.get())
   directory = filedialog.askdirectory()
   win.quit()

#Initialize a Label to display the User Input
label=tk.Label(win, text="Input tanggal, bulan dan tahun bacaan yang diinginkan", font=("Courier 15 bold"))
label.grid(row=0, column=0, columnspan=3)

#Create Entry widgets to accept User Input
L1 = tk.Label(text="Tanggal")
L1.grid(row=1,column=0)

entry_tanggal= tk.Entry(win)
entry_tanggal.focus_set()
entry_tanggal.insert(0, f"{today.day}")
entry_tanggal.grid(row=1,column=1)

L2 = tk.Label(text="Bulan")
L2.grid(row=2, column=0)

entry_bulan= tk.Entry(win)
entry_bulan.insert(0, f"{today.month}")
entry_bulan.grid(row=2, column=1)

L3 = tk.Label(text="Tahun")
L3.grid(row=3,column=0)

entry_tahun = tk.Entry(win)
entry_tahun.insert(0, f"{today.year}")
entry_tahun.grid(row=3,column=1)

L4 = tk.Label(text="Fontsize")
L4.grid(row=4,column=0)

entry_fontsize = tk.Entry(win)
entry_fontsize.grid(row=4,column=1)
entry_fontsize.insert(0, "28")

L5 = tk.Label(text="Maximum Karakter(termasuk spasi) dalam slide")
L5.grid(row=5,column=0)


entry_wordcount = tk.Entry(win)
entry_wordcount.grid(row=5,column=1)
entry_wordcount.insert(0, "500")

#Create a Button to validate Entry Widget
ttk.Button(win, text= "Confirm",width= 20, command= get_text).grid(column=1, row=6)

win.mainloop()

url=f"http://www.imankatolik.or.id/kalender.php?b={bulan}&t={tahun}"
base_url= "http://www.imankatolik.or.id"

#Website Kalender
r=requests.get(url)
html_doc=r.text
soup = BeautifulSoup(html_doc, features="lxml")
table = soup.find_all("td", class_="k_alkitab")

links=table[int(tanggal)-1].find_all("a")
links.pop()

prs= Presentation("Default.pptx")
slide1= prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Powerpoint Misa"
subtitle.text = f"{tanggal}/{bulan}/{tahun}"

#Website Ayat
for link in links:
    ayat=requests.get(base_url + link.get("href"))
    ayat_html= BeautifulSoup(ayat.text, features="lxml")
    tr=ayat_html.find_all("tr")
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    shapes = slide2.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text= "; ".join(link.get("href")[16:].split(";"))
    tf = body_shape.text_frame

    counter=0
    word_counter= 0

    for td in tr:
        kalimat= list(list(td.children)[1].stripped_strings)[0]

        if counter==0:
            tf.text=kalimat
            p=tf.add_paragraph()
            counter+=1
            word_counter += len(kalimat)
        else:
            if (word_counter + len(kalimat) > word_limit):
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(fontsize)

                slide2 = prs.slides.add_slide(prs.slide_layouts[1])
                shapes = slide2.shapes

                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                title_shape.text= "; ".join(link.get("href")[16:].split(";"))
                tf = body_shape.text_frame

                word_counter = len(kalimat)
                tf.text=kalimat
                p=tf.add_paragraph()
            else:
                p.text=kalimat
                p=tf.add_paragraph()
                word_counter += len(kalimat)

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(fontsize)

prs.save(r"{}/PPT_Misa.pptx".format(directory))
